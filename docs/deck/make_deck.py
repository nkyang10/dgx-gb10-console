#!/usr/bin/env python3
"""Generate client-facing deck for the DGX GB10 Console product.
python-pptx deck, warm-paper theme (user design DNA): #f7f4ef bg, #3e5c4b ink-green,
#c05f3c terracotta accent, #2f3e38 dark text, Noto Serif TC headings.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- theme ----
PAPER   = RGBColor(0xF7, 0xF4, 0xEF)
INK     = RGBColor(0x3E, 0x5C, 0x4B)   # ink green
CLAY    = RGBColor(0xC0, 0x5F, 0x3C)   # terracotta
DARK    = RGBColor(0x2E, 0x3A, 0x34)   # text
MUTED   = RGBColor(0x6B, 0x74, 0x6E)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LINE    = RGBColor(0xE4, 0xDF, 0xD5)

FONT_T   = "Noto Serif TC"   # headings
FONT_B   = "Noto Sans TC"    # body

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def _set_fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color

def _no_line(shape):
    shape.line.fill.background()

def _font(run, size, bold=False, color=DARK, name=FONT_B):
    f = run.font; f.size = Pt(size); f.bold = bold; f.name = name
    f.color.rgb = color
    # ensure CJK serif applies via latin too
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea')); 
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', name)

def add_slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH); _set_fill(r, bg); _no_line(r)
    r.shadow.inherit = False
    return s

def textbox(s, x, y, w, h):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    return tf

def para(tf, text, size, bold=False, color=DARK, name=FONT_T, align=PP_ALIGN.LEFT,
         space_before=0, space_after=6, first=False, line=None):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before); p.space_after = Pt(space_after)
    if line: p.line_spacing = line
    r = p.add_run(); r.text = text
    _font(r, size, bold, color, name)
    return p

def rect(s, x, y, w, h, color, line_w=0):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _set_fill(shp, color); shp.line.fill.background(); shp.shadow.inherit=False
    return shp

def chip(s, x, y, w, text, color=INK, bg=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.34))
    if bg: _set_fill(shp, bg)
    else:  _set_fill(shp, color)
    shp.line.color.rgb = color; shp.line.width = Pt(1); shp.shadow.inherit=False
    tf = shp.text_frame; tf.word_wrap=False
    tf.margin_left=Inches(0.06); tf.margin_right=Inches(0.06); tf.margin_top=0; tf.margin_bottom=0
    p = tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=text; _font(r, 10, bold=True, color=(WHITE if bg else color), name=FONT_B)
    return shp

def footer(s, idx, total):
    tf = textbox(s, Inches(0.7), SH-Inches(0.42), SW-Inches(1.4), Inches(0.3))
    para(tf, f"DGX GB10 Console  ·  {idx}/{total}", 9, color=MUTED, name=FONT_B)

TOTAL = 14
def bullet(tf, text, size=15, color=DARK, mark="— ", bold_mark=False):
    p = tf.add_paragraph(); p.space_after=Pt(5); p.line_spacing=1.12
    r0=p.add_run(); r0.text=mark; _font(r0, size, bold=True, color=CLAY, name=FONT_B)
    r=p.add_run(); r.text=text; _font(r, size, bold=bold_mark, color=color, name=FONT_B)

# ============ SLIDE 1: Cover ============
s = add_slide(INK)
rect(s, 0, 0, SW, Inches(0.14), CLAY)
tf = textbox(s, Inches(1.0), Inches(2.35), SW-Inches(2), Inches(2.6))
para(tf, "DGX GB10 Console", 60, bold=True, color=WHITE, name=FONT_T, first=True, space_after=10)
para(tf, "DGX Spark 集群 即時狀態 × 效能 × 健康監控平台", 22, color=RGBColor(0xD9,0xE4,0xDD), name=FONT_B, space_after=28)
para(tf, "一部 HTTP API，實時睇晒你部 DGX 做緊咩 — 硬體、軟體、vLLM 推理效能、成本、告警", 15, color=RGBColor(0xC7,0xD4,0xCC), name=FONT_B)
tf2 = textbox(s, Inches(1.0), SH-Inches(1.15), SW-Inches(2), Inches(0.6))
para(tf2, "適用於 2 至 N 部 DGX Spark（GB10 Grace Blackwell）", 13, color=RGBColor(0xA9,0xB9,0xB0), name=FONT_B, first=True)
footer(s,1,TOTAL)

# ============ SLIDE 2: Pain / Value ============
s = add_slide()
rect(s, 0, 0, Inches(0.14), SH, CLAY)
tf = textbox(s, Inches(1.0), Inches(0.6), SW-Inches(2), Inches(0.9))
para(tf, "客戶痛點：你很哩緊部 DGX，但「佢喺度做緊乜？」", 30, bold=True, color=INK, first=True, space_after=8)
para(tf, "無可視化 → 過熱降頻、OOM、排隊爆煲、磁碟滿、成本失控，全部唔知", 15, color=MUTED, name=FONT_B)
tf = textbox(s, Inches(1.0), Inches(2.0), SW-Inches(2), Inches(3.6))
bullet(tf, "只知道部機「著咗」，但唔知 GPU 係咪過熱緊→降頻，效能偷偷流失")
bullet(tf, "vLLM 排緊幾多 request？TTFT 有冇爆炸？邊個 model 最夾樽頸？")
bullet(tf, "磁碟幾時滿？NVMe 幾時壞？能耗同 Token 成本每月幾多？")
bullet(tf, "要逐部機 SSH 去睇 nvidia-smi，N 部機睇到頭都大")
bullet(tf, "出事先發現，SLA／收入無形損失")
footer(s,2,TOTAL)

# ============ SLIDE 3: Product overview ============
s = add_slide()
tf = textbox(s, Inches(1.0), Inches(0.55), SW-Inches(2), Inches(0.8))
para(tf, "產品：一部 HTTP API 睇晒成個集群", 30, bold=True, color=INK, first=True, space_after=6)
para(tf, "N 節點星型架構 · 每部 DGX 一個 exporter · 一部 aggregator 合併", 14, color=MUTED, name=FONT_B)
# 4 feature cards
cards = [
    ("硬體健康", "GPU / CPU / 記憶體(UMA) / 溫度 / 功率\n過熱降頻矩陣·節流偵測 · NVMe 壽命"),
    ("vLLM 效能", "tok/s · TTFT / ITL / e2e · queue\nKV cache · per-model · SLO"), 
    ("成本 / 容量", "能耗與電費 · Token 會計\n磁碟 TTF · KV / 記憶體 headroom"),
    ("可靠 / 安全", "container OOM · kernel log\nsecurity audit · RDMA 擁塞"),
]
cw = Inches(2.82); gap = Inches(0.24); x0 = Inches(1.0); y0 = Inches(1.85); ch = Inches(3.0)
for i,(t,d) in enumerate(cards):
    x = x0 + i*(cw+gap)
    card = rect(s, x, y0, cw, ch, RGBColor(0xF0,0xEC,0xE3)); card.line.color.rgb=LINE; card.line.width=Pt(1)
    rect(s, x, y0, cw, Inches(0.10), CLAY if i%2==0 else INK)
    tf2 = textbox(s, x+Inches(0.3), y0+Inches(0.4), cw-Inches(0.6), ch-Inches(0.6))
    para(tf2, t, 19, bold=True, color=INK, first=True, space_after=8)
    for ln in d.split("\n"):
        para(tf2, ln, 13, color=DARK, name=FONT_B, space_after=3)
footer(s,3,TOTAL)

# ============ SLIDE 4: Scale ============
s = add_slide()
tf = textbox(s, Inches(1.0), Inches(0.6), SW-Inches(2), Inches(0.8))
para(tf, "由 2 部到 N 部，輕鬆擴展", 30, bold=True, color=INK, first=True, space_after=6)
para(tf, "星型架構 · 每部節點都係一等公民 · 加機只需 3 步", 14, color=MUTED, name=FONT_B)
# diagram: aggregator center hub + 3 leaves
rect(s, Inches(4.9), Inches(2.1), Inches(3.5), Inches(1.2), INK)
tf2=textbox(s, Inches(4.9), Inches(2.45), Inches(3.5), Inches(0.6)); para(tf2,"Aggregator (dgx-01)", 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, first=True)
rect(s, Inches(1.6), Inches(4.4), Inches(3.0), Inches(1.0), RGBColor(0xE9,0xE4,0xD8))
rect(s, Inches(5.2), Inches(4.4), Inches(3.0), Inches(1.0), RGBColor(0xE9,0xE4,0xD8))
rect(s, Inches(8.8), Inches(4.4), Inches(3.0), Inches(1.0), RGBColor(0xE9,0xE4,0xD8))
for (x,i) in [(1.9,2),(5.5,3),(9.1,4)]:
    tf3=textbox(s, Inches(x), Inches(4.7), Inches(2.4), Inches(0.4)); para(tf3,f"dgx-{i:02d} exporter", 13, bold=True, color=INK, align=PP_ALIGN.CENTER, first=True)
tf4=textbox(s, Inches(1.0), Inches(6.0), SW-Inches(2), Inches(0.5))
para(tf4,"add 一部：改 config 加一行 nodes[] + 喺新機跑 install.sh --node-id=<id> —— 自動出現喺 snapshot", 13, color=MUTED, align=PP_ALIGN.CENTER, name=FONT_B, first=True)
footer(s,4,TOTAL)

# ============ SLIDE 5: Hardware health + thermal ============
s = add_slide()
tf = textbox(s, Inches(1.0), Inches(0.55), SW-Inches(2), Inches(0.8))
para(tf, "硬體健康 · 過熱降頻一眼睇到", 30, bold=True, color=INK, first=True, space_after=6)
para(tf, "GB10 統一記憶體(UMA) 專屬做法 — 唔靠報 N/A 嘅 nvidia-smi VRAM", 14, color=MUTED, name=FONT_B)
tf2 = textbox(s, Inches(1.0), Inches(1.8), Inches(6.6), Inches(4.6))
bullet(tf2, "GPU：util / 溫度 / 功率 / 時脈 / throttle reason（熱 vs 功率 vs 卡死）")
bullet(tf2, "熱降頻觀測矩陣：溫度 + SM 頻率比率 + throttle bit + 負載 夾埋判")
bullet(tf2, "CPU / 記憶體(UMA) / 磁碟(NVMe SMART 壽命) / 網絡(Fabric RDMA)")
bullet(tf2, "節流史 timeline：幾時降頻、持續幾耐 — 自動計入告警")
bullet(tf2, "GB10 特有「卡死低時脈」bug 偵測（USB-C PD 供電）")
tf3 = textbox(s, Inches(8.0), Inches(1.8), Inches(4.3), Inches(4.6))
tc = rect(s, Inches(8.0), Inches(1.85), Inches(4.3), Inches(2.9), RGBColor(0xF0,0xEC,0xE3)); tc.line.color.rgb=LINE
para(tf3, "熱降頻判讀", 16, bold=True, color=INK, first=True, space_after=6)
for a,b in [("temp>80 + clock 低 + thermal bit","過熱降頻 🔴"), ("power_cap bit","功率上限 🟠"), ("clock<45% + 無 reason","卡死 bug 🔴")]:
    para(tf3, f"{a}", 12, color=DARK, name=FONT_B, space_after=1)
    para(tf3, f"      →  {b}", 12, bold=True, color=CLAY, name=FONT_B, space_after=6)
footer(s,5,TOTAL)

# ============ SLIDE 6: vLLM performance ============
s = add_slide()
tf = textbox(s, Inches(1.0), Inches(0.55), SW-Inches(2), Inches(0.8))
para(tf, "vLLM 推理效能 — 個 API 真正主角", 30, bold=True, color=INK, first=True, space_after=6)
para(tf, "喺 Docker 入面跑嘅 vLLM，直接撈 /metrics · /health · /v1/models · /is_sleeping", 14, color=MUTED, name=FONT_B)
tf2 = textbox(s, Inches(1.0), Inches(1.8), Inches(6.6), Inches(4.6))
bullet(tf2, "吞吐：decode / prefill tok/s，RPS（完成 request 率）")
bullet(tf2, "延遲：TTFT、ITL/TPOT、e2e — avg + P95（窗口度率，唔抄即時數）")
bullet(tf2, "積壓：running / waiting 隊列，per-model backlog，queue time")
bullet(tf2, "KV-cache 用率、prefix-cache hit rate、preemption")
bullet(tf2, "per-model SLO goodput；engine sleep / idle / dead 狀態")
tf3 = textbox(s, Inches(8.0), Inches(1.8), Inches(4.3), Inches(4.0))
tc = rect(s, Inches(8.0), Inches(1.85), Inches(4.3), Inches(3.4), RGBColor(0xF0,0xEC,0xE3)); tc.line.color.rgb=LINE
tf3b = textbox(s, Inches(8.4), Inches(2.1), Inches(3.5), Inches(3.0))
para(tf3b, "想知佢有冇排隊？", 16, bold=True, color=INK, first=True, space_after=6)
para(tf3b, "num_requests_waiting ↑  +  queue_time ↑  = 積壓緊、處理緊", 13, color=DARK, name=FONT_B, space_after=10)
para(tf3b, "唔係掛咗 — 係 vLLM 自己排隊消化緊", 13, bold=True, color=CLAY, name=FONT_B)
footer(s,6,TOTAL)

# ============ SLIDE 7: Cost / capacity ============
s = add_slide()
tf = textbox(s, Inches(1.0), Inches(0.55), SW-Inches(2), Inches(0.8))
para(tf, "成本透明 · 容量先知", 30, bold=True, color=INK, first=True, space_after=6)
para(tf, "唔使估 — 幾多電、幾多 Token、幾時滿，全部數字化", 14, color=MUTED, name=FONT_B)
tf2 = textbox(s, Inches(1.0), Inches(1.8), Inches(6.6), Inches(4.6))
bullet(tf2, "能耗：整機 wall power（spbm），kWh/日，電費（$/kWh）")
bullet(tf2, "idle vs active 成本分帳：幾多電係白拎")
bullet(tf2, "Token 會計：per-model 每日 tokens、$/1M tokens、cache 慳幾多")
bullet(tf2, "預算告警：$ 突破 80% / 100% 即知")
header2 = textbox(s, Inches(8.0), Inches(1.8), Inches(4.3), Inches(4.6))
para(header2, "容量預測", 18, bold=True, color=INK, first=True, space_after=6)
tc = rect(s, Inches(8.0), Inches(2.3), Inches(4.3), Inches(3.0), RGBColor(0xF0,0xEC,0xE3)); tc.line.color.rgb=LINE
tf3 = textbox(s, Inches(8.4), Inches(2.6), Inches(3.5), Inches(2.4))
bullet(tf3, "disk TTF：幾時 full（4h / 6h 預警）", 14)
bullet(tf3, "KV / 記憶體 headroom 趨勢（UMA）", 14)
bullet(tf3, "「要唔要加 node」呢啲問題有數據答", 14)
footer(s,7,TOTAL)

# ============ SLIDE 8: Reliability & security ============
s = add_slide()
tf = textbox(s, Inches(1.0), Inches(0.55), SW-Inches(2), Inches(0.8))
para(tf, "可靠監察 · 安全守門", 30, bold=True, color=INK, first=True, space_after=6)
tf2 = textbox(s, Inches(1.0), Inches(1.75), Inches(6.6), Inches(5.0))
bullet(tf2, "NVMe 壽命預測 + critical warning（呢個平台有已知 No-NVMe 故障）")
bullet(tf2, "container OOM / restart loop → 即時警報")
bullet(tf2, "kernel / OOM / NVRM 日誌 markers（24h 有事冇事一目了然）")
bullet(tf2, "睇 log：`/logs` 攞 vLLM / kernel / 服務 / container 日誌（限行數 + auth）")
bullet(tf2, "security audit：open ports、failed SSH、異常 process/container")
tf3 = textbox(s, Inches(8.0), Inches(1.75), Inches(4.3), Inches(5.0))
tc=rect(s, Inches(8.0), Inches(1.8), Inches(4.3), Inches(3.4), RGBColor(0xF0,0xEC,0xE3)); tc.line.color.rgb=LINE
tf3b=textbox(s, Inches(8.4), Inches(2.1), Inches(3.5), Inches(2.8))
para(tf3b,"安全",17,bold=True,color=INK,first=True,space_after=6)
bullet(tf3b,"API read-only 預設 + Token / API-Key",14)
bullet(tf3b,"站與站用 mTLS 或 WireGuard 加密",14)
bullet(tf3b,"vLLM 控制（可選）受控 + 唔 auto-wake",14)
footer(s,8,TOTAL)

# ============ SLIDE 9: Alerts ============
s = add_slide()
tf = textbox(s, Inches(1.0), Inches(0.55), SW-Inches(2), Inches(0.8))
para(tf, "即時告警 — 未出事已經話你知", 30, bold=True, color=INK, first=True, space_after=6)
para(tf, "成個集群一套 alert 規則，Predomtheus / Grafana 或內建 /metrics", 14, color=MUTED, name=FONT_B)
rows = [
    ("NVMe wear / critical warning", "將壞碟，先換先安","crit"),
    ("GPU > 96°C / SM clock 崩", "過熱降頻 / USB-C PD 卡死","crit"),
    ("Disk TTF < 6h", "磁碟就快快滿","warn"),
    ("KV cache 飽和 + 排隊 + preempt", "模型唔夠位，要加 node","warn"),
    ("RDMA link down / 擁塞升", "fabric 出問題","crit"),
    ("vLLM engine dead (health≠200)","推理死機，即刻喊","crit"),
]
y0=Inches(1.8); rh=Inches(0.62)
for i,(a,b,sev) in enumerate(rows):
    y=y0+i*rh
    wide = rect(s, Inches(1.0), y, Inches(8.6), rh-Inches(0.12), RGBColor(0xF0,0xEC,0xE3)); wide.line.color.rgb=LINE; wide.line.width=Pt(0.75)
    tfa=textbox(s, Inches(1.25), y+Inches(0.06), Inches(5.0), rh-Inches(0.2)); para(tfa,a,14,bold=True,color=DARK,first=True)
    tfb=textbox(s, Inches(6.35), y+Inches(0.06), Inches(3.1), rh-Inches(0.2)); para(tfb,b,12,color=MUTED,name=FONT_B,first=True)
    chip(s, Inches(9.8), y+Inches(0.12), Inches(1.5), sev.upper(), color=(CLAY if sev=="crit" else INK), bg=RGBColor(0xF2,0x6E,0x4A) if sev=="crit" else False)
footer(s,9,TOTAL)

# ============ SLIDE 10: Architecture ============
s = add_slide()
tf = textbox(s, Inches(1.0), Inches(0.55), SW-Inches(2), Inches(0.8))
para(tf, "架構 — 分層清晰，通訊安全", 30, bold=True, color=INK, first=True, space_after=6)
tf2 = textbox(s, Inches(1.0), Inches(1.7), Inches(6.0), Inches(4.8))
bullet(tf2, "L4 Consumer：你部機 / Prometheus / Grafana")
bullet(tf2, "L3 API：aggregator（/health · /cluster.json · /metrics · /logs · /control") 
bullet(tf2, "L2 Exporter：每部 DGX 一個 :9101（一等公民）")
bullet(tf2, "L1 Collectors：hardware · system · network · services · vllm（並行）")
bullet(tf2, "L0 數據源：nvidia-smi / proc / sys / docker / vLLM / RDMA")
tf3 = textbox(s, Inches(7.4), Inches(1.7), Inches(5.0), Inches(4.8))
tc=rect(s, Inches(7.4), Inches(1.75), Inches(5.0), Inches(3.2), RGBColor(0xF0,0xEC,0xE3)); tc.line.color.rgb=LINE
tf3b=textbox(s, Inches(7.8), Inches(2.1), Inches(4.2), Inches(2.5))
para(tf3b,"通訊安全",17,bold=True,color=INK,first=True,space_after=6)
bullet(tf3b,"出外：token + read-only + TLS",14)
bullet(tf3b,"站與站：mTLS / WireGuard 加密",14)
bullet(tf3b,"/logs 敏感 → auth；/control 受控",14)
footer(s,10,TOTAL)

# ============ SLIDE 11: Deployment ============
s = add_slide()
tf = textbox(s, Inches(1.0), Inches(0.55), SW-Inches(2), Inches(0.8))
para(tf, "部署 — 幾分鐘上線、遷移就把炮", 30, bold=True, color=INK, first=True, space_after=6)
tf2 = textbox(s, Inches(1.0), Inches(1.8), Inches(6.6), Inches(4.6))
bullet(tf2, "每節點 venv + systemd（原裝 DGX OS，ARM64 wheels 齊）")
bullet(tf2, "冪等 install.sh：`--role --node-id --mgmt-ip`，加機 3 步")
bullet(tf2, "hardened systemd unit（NoNewPrivileges / ProtectSystem …）")
bullet(tf2, "可選 Docker 路徑")
bullet(tf2, "換機／升級：`install.sh --update`／`--uninstall` 回滾")
col2 = textbox(s, Inches(8.0), Inches(1.8), Inches(4.3), Inches(4.6))
tc=rect(s, Inches(8.0), Inches(1.85), Inches(4.3), Inches(2.6), RGBColor(0xF0,0xEC,0xE3)); tc.line.color.rgb=LINE
tf3=textbox(s, Inches(8.4), Inches(2.2), Inches(3.5), Inches(2.0))
para(tf3,"部署三句",17,bold=True,color=INK,first=True,space_after=8)
para(tf3,"① clone + install.sh",14,color=DARK,name=FONT_B,space_after=6)
para(tf3,"② 定 token / 角色",14,color=DARK,name=FONT_B,space_after=6)
para(tf3,"③ 起 service — 完成",14,color=DARK,name=FONT_B)
footer(s,11,TOTAL)

# ============ SLIDE 12: Integration ============
s = add_slide()
tf = textbox(s, Inches(1.0), Inches(0.55), SW-Inches(2), Inches(0.8))
para(tf, "整合 — 唔係孤島", 30, bold=True, color=INK, first=True, space_after=6)
tf2 = textbox(s, Inches(1.0), Inches(1.8), Inches(6.6), Inches(4.6))
bullet(tf2, "Prometheus / Grafana：內建 /metrics 直接被撈 → dashboard + alert")
bullet(tf2, "Webhook / Telegram（可選）：出事即 push")
bullet(tf2, "SQLite 歷史 / ring buffer：即時 + 回溯（可配 Prometheus TSDB 長史）")
bullet(tf2, "vLLM 控制（可選）：abort-by-id、pause/resume — 受控 + 唔干預")
tf3 = textbox(s, Inches(8.0), Inches(1.8), Inches(4.3), Inches(4.6))
tc=rect(s, Inches(8.0), Inches(1.85), Inches(4.3), Inches(2.2), RGBColor(0xF0,0xEC,0xE3)); tc.line.color.rgb=LINE
tf3b=textbox(s, Inches(8.4), Inches(2.2), Inches(3.5), Inches(1.6))
para(tf3b,"想擴到大型集群？",16,bold=True,color=INK,first=True,space_after=6)
para(tf3b,"N 前預設單 aggregator；N 大到可 federation / 多級",13,color=DARK,name=FONT_B)
footer(s,12,TOTAL)

# ============ SLIDE 13: Value summary ============
s = add_slide(INK)
rect(s, 0, 0, SW, Inches(0.14), CLAY)
tf = textbox(s, Inches(1.0), Inches(1.0), SW-Inches(2), Inches(5.2))
para(tf, "客戶價值", 40, bold=True, color=WHITE, first=True, space_after=18)
for mark,txt in [("✓","唔使估 — 成個集群即時狀態、效能、成本一目了然"),
                 ("✓","過熱降頻 / OOM / 排隊 / 磁碟滿 / 將壞碟 — 未出事已經話你知"),
                 ("✓","效能 SLA 有得查：TTFT / P95 / per-model — 唔使靠估"),
                 ("✓","成本透明：電費、Token、容量預測 — 慳錢決定有數據"),
                 ("✓","安全守門：read-only + mTLS + 受控操作")]:
    p=tf.add_paragraph(); p.space_after=Pt(14)
    r0=p.add_run(); r0.text=mark+"  "; _font(r0,22,bold=True,color=CLAY,name=FONT_B)
    r=p.add_run(); r.text=txt; _font(r,20,color=WHITE,name=FONT_B)
tf2=textbox(s, Inches(1.0), SH-Inches(1.1), SW-Inches(2), Inches(0.6))
para(tf2,"由「唔知」到「一個 API 就知」— DGX GB10 Console", 14, color=RGBColor(0xC7,0xD4,0xCC), name=FONT_B, first=True)
footer(s,13,TOTAL)

# ============ SLIDE 14: Next steps / close ============
s = add_slide(PAPER)
tf = textbox(s, Inches(1.0), Inches(2.1), SW-Inches(2), Inches(3.0))
para(tf, "下一 步", 40, bold=True, color=INK, first=True, space_after=16)
para(tf, "落地試行：喺你嘅 DGX Spark 上部署概念驗證", 20, color=DARK, name=FONT_B, space_after=8)
para(tf, "接入你嘅 vLLM / 監控棧，睇住真實數據", 20, color=DARK, name=FONT_B, space_after=8)
para(tf, "按你需要收窄或擴展收嘅指標", 20, color=DARK, name=FONT_B, space_after=8)
tf2 = textbox(s, Inches(1.0), Inches(6.0), SW-Inches(2), Inches(0.8))
para(tf2, "DGX GB10 Console — 謝謝", 18, bold=True, color=CLAY, first=True)
footer(s,14,TOTAL)

prs.save("/root/projects/dgx-status-api/docs/deck/DGX_GB10_Console.pptx")
print("saved pptx")
